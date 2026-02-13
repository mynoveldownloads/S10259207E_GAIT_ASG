import os
import base64
import io
import fitz  # pymupdf
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
try:
    from docx import Document
except ImportError:
    Document = None
from openai import OpenAI

# Initialize OpenRouter Client
# In a production environment, use environment variables. 
# Using the provided key as requested.
API_KEY = "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"
BASE_URL = "https://openrouter.ai/api/v1"
MODEL_NAME = "mistralai/ministral-8b-2512"

client = OpenAI(
  base_url=BASE_URL,
  api_key=API_KEY,
)

def encode_bytes_to_base64(image_bytes: bytes, mime_type: str = "image/png") -> str:
    """Encodes image bytes to a data URL for the API."""
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    return f"data:{mime_type};base64,{base64_image}"

def transcribe_image(image_bytes: bytes, mime_type: str = "image/png") -> str:
    """
    Transcribes text from a single image using the Ministral model.
    """
    image_url = encode_bytes_to_base64(image_bytes, mime_type)
    
    messages = [
        {
            "role": "system",
            "content": "You are an expert OCR assistant. Your task is to extract all visible text from the image accurately. If it's a technical diagram, describe the labels and relationships briefly."
        },
        {
            "role": "user",
            "content": [
                {"type": "text", "text": "Transcribe the text in this image."},
                {"type": "image_url", "image_url": {"url": image_url}}
            ]
        }
    ]
    
    try:
        completion = client.chat.completions.create(
            extra_headers={
                "HTTP-Referer": "http://localhost:8501",
                "X-Title": "Streamlit Document Parser",
            },
            model=MODEL_NAME,
            messages=messages
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error transcribing image: {e}")
        return f"[OCR Error: {e}]"

def process_pdf_content(file_path: str):
    """
    Extracts plaintext and transcribes images from each page of the PDF.
    Returns a string containing all extracted text structured by page.
    """
    full_text = []
    doc = fitz.open(file_path)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # 1. Try to extract plaintext directly
        page_text = page.get_text().strip()
        
        # 2. Check for images on the page
        image_list = page.get_images(full=True)
        
        extracted_content = [f"--- PDF Page {page_num + 1} ---"]
        
        if page_text:
            extracted_content.append(page_text)
        
        # 3. If there are images, render the page as an image for transcription
        if image_list:
            pix = page.get_pixmap()
            img_bytes = pix.tobytes("png")
            
            transcribed_text = transcribe_image(img_bytes, "image/png")
            extracted_content.append(f"[Image Content Analysis]:\n{transcribed_text}")
            
        full_text.append("\n".join(extracted_content))
        
    doc.close()
    return "\n\n".join(full_text)

def process_pptx_content(file_path: str):
    """
    Extracts text and transcribes images from PPTX slides.
    Returns a string containing all extracted text structured by slide.
    """
    full_text = []
    prs = Presentation(file_path)
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_image_transcriptions = []
        
        # Extract Text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Extract and transcribe Images individually
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob
                    ext = shape.image.ext
                    mime_type = f"image/{ext}" if ext else "image/png"
                    
                    transcribed_text = transcribe_image(image_blob, mime_type)
                    slide_image_transcriptions.append(transcribed_text)
                except Exception as e:
                    print(f"Error extracting image from slide {i+1}: {e}")

        slide_content = [f"--- Slide {i + 1} ---"]
        if slide_text:
            slide_content.append("\n".join(slide_text))
        
        if slide_image_transcriptions:
            slide_content.append("[Image Content Analysis]:")
            slide_content.extend(slide_image_transcriptions)
            
        full_text.append("\n".join(slide_content))
            
    return "\n\n".join(full_text)

def process_docx_content(file_path: str):
    """
    Extracts text from DOCX files.
    """
    if Document is None:
        return "Error: 'python-docx' is not installed. Please run 'pip install python-docx' in your environment."
    
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())
    return "\n".join(full_text)

def get_ocr_content(file_path: str) -> str:
    """
    Processes a PDF or PPTX file and returns the aggregated text and OCR transcriptions.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            extracted_text = process_pdf_content(file_path)
        elif ext == ".pptx":
            extracted_text = process_pptx_content(file_path)
        elif ext == ".docx":
            extracted_text = process_docx_content(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            with open(file_path, "rb") as f:
                img_bytes = f.read()
            mime_type = f"image/{ext[1:]}"
            extracted_text = transcribe_image(img_bytes, mime_type)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
            
        if not extracted_text:
            return "Error: No content could be extracted from the file."

        return extracted_text
        
    except Exception as e:
        return f"Error in OCR Processing: {str(e)}"
